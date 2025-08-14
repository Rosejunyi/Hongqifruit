#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
PPT内容提取工具
提取PowerPoint文件中的文本和图片内容
"""

import os
import sys
import json
from pathlib import Path

def extract_ppt_content_basic(ppt_path):
    """
    基础方法：尝试从PPT二进制文件中提取可读文本
    """
    try:
        with open(ppt_path, 'rb') as file:
            content = file.read()
            
        # 尝试解码为UTF-8文本，忽略错误
        try:
            text_content = content.decode('utf-8', errors='ignore')
        except:
            text_content = content.decode('gbk', errors='ignore')
            
        # 提取可能的文本片段（长度大于3的连续字符）
        import re
        
        # 查找中文文本
        chinese_pattern = r'[\u4e00-\u9fff]+'
        chinese_texts = re.findall(chinese_pattern, text_content)
        chinese_texts = [text for text in chinese_texts if len(text) >= 2]
        
        # 查找英文文本
        english_pattern = r'[A-Za-z][A-Za-z\s]{3,50}'
        english_texts = re.findall(english_pattern, text_content)
        english_texts = [text.strip() for text in english_texts if len(text.strip()) >= 4]
        
        # 查找数字和日期
        number_pattern = r'\d{4}[-年]\d{1,2}[-月]\d{1,2}[日]?|\d{4}年|\d+%|\d+万|\d+千|\d+元'
        numbers = re.findall(number_pattern, text_content)
        
        # 查找可能的联系方式
        contact_pattern = r'1[3-9]\d{9}|[\w\.-]+@[\w\.-]+\.\w+|www\.[\w\.-]+|http[s]?://[\w\.-]+'
        contacts = re.findall(contact_pattern, text_content)
        
        return {
            'chinese_texts': list(set(chinese_texts)),
            'english_texts': list(set(english_texts)),
            'numbers_dates': list(set(numbers)),
            'contacts': list(set(contacts)),
            'file_size': len(content)
        }
        
    except Exception as e:
        return {'error': str(e)}

def extract_with_python_pptx():
    """
    尝试使用python-pptx库提取内容（如果可用）
    """
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        
        ppt_path = "../红琪食品公司简介20240305.pptx"
        
        if not os.path.exists(ppt_path):
            return {'error': 'PPT文件不存在'}
            
        prs = Presentation(ppt_path)
        
        extracted_data = {
            'slides': [],
            'total_slides': len(prs.slides),
            'images': []
        }
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_data = {
                'slide_number': slide_idx + 1,
                'texts': [],
                'titles': [],
                'images_count': 0
            }
            
            for shape in slide.shapes:
                # 提取文本
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if len(text) > 1:
                        if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                            # 检查是否为标题
                            first_para = shape.text_frame.paragraphs[0]
                            if hasattr(first_para, 'font') and first_para.font.size:
                                if first_para.font.size.pt > 18:  # 大字体可能是标题
                                    slide_data['titles'].append(text)
                                else:
                                    slide_data['texts'].append(text)
                            else:
                                slide_data['texts'].append(text)
                        else:
                            slide_data['texts'].append(text)
                
                # 统计图片
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    slide_data['images_count'] += 1
                    
            extracted_data['slides'].append(slide_data)
            
        return extracted_data
        
    except ImportError:
        return {'error': 'python-pptx库未安装，请运行: pip install python-pptx'}
    except Exception as e:
        return {'error': f'提取失败: {str(e)}'}

def main():
    """主函数"""
    print("=== PPT内容提取工具 ===")
    
    ppt_path = "../红琪食品公司简介20240305.pptx"
    
    if not os.path.exists(ppt_path):
        print(f"错误：找不到文件 {ppt_path}")
        return
    
    print(f"正在处理文件: {ppt_path}")
    print(f"文件大小: {os.path.getsize(ppt_path) / 1024:.1f} KB")
    
    # 方法1：基础二进制提取
    print("\n--- 方法1：基础文本提取 ---")
    basic_result = extract_ppt_content_basic(ppt_path)
    
    if 'error' in basic_result:
        print(f"基础提取失败: {basic_result['error']}")
    else:
        print(f"找到中文文本片段: {len(basic_result['chinese_texts'])}个")
        print("中文内容示例:")
        for text in basic_result['chinese_texts'][:10]:
            if len(text) >= 3:
                print(f"  - {text}")
        
        print(f"\n找到英文文本片段: {len(basic_result['english_texts'])}个")
        print("英文内容示例:")
        for text in basic_result['english_texts'][:5]:
            print(f"  - {text}")
            
        if basic_result['numbers_dates']:
            print(f"\n找到数字/日期: {len(basic_result['numbers_dates'])}个")
            for num in basic_result['numbers_dates'][:10]:
                print(f"  - {num}")
                
        if basic_result['contacts']:
            print(f"\n找到联系方式: {len(basic_result['contacts'])}个")
            for contact in basic_result['contacts']:
                print(f"  - {contact}")
    
    # 方法2：使用python-pptx库
    print("\n--- 方法2：专业PPT解析 ---")
    pptx_result = extract_with_python_pptx()
    
    if 'error' in pptx_result:
        print(f"专业解析失败: {pptx_result['error']}")
        print("建议安装python-pptx: pip install python-pptx")
    else:
        print(f"成功解析PPT，共 {pptx_result['total_slides']} 张幻灯片")
        
        for slide in pptx_result['slides']:
            print(f"\n第 {slide['slide_number']} 张幻灯片:")
            if slide['titles']:
                print(f"  标题: {', '.join(slide['titles'])}")
            if slide['texts']:
                print(f"  文本内容: {len(slide['texts'])} 段")
                for text in slide['texts'][:3]:  # 只显示前3段
                    print(f"    - {text[:50]}...")
            if slide['images_count']:
                print(f"  图片数量: {slide['images_count']}")
    
    # 保存结果到JSON文件
    output_file = "ppt_extracted_content.json"
    combined_result = {
        'basic_extraction': basic_result,
        'pptx_extraction': pptx_result,
        'extraction_time': str(Path().cwd()),
        'source_file': ppt_path
    }
    
    try:
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(combined_result, f, ensure_ascii=False, indent=2)
        print(f"\n结果已保存到: {output_file}")
    except Exception as e:
        print(f"保存结果失败: {e}")

if __name__ == "__main__":
    main()
